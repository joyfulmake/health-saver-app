from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── helpers ──────────────────────────────────────────────────────────────────
W, H = Inches(13.33), Inches(7.5)   # 16:9 widescreen

BG_DARK   = RGBColor(0x0f, 0x0f, 0x0f)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MUTED     = RGBColor(0x88, 0x88, 0x88)
GREEN     = RGBColor(0x16, 0xa3, 0x4a)
BLUE      = RGBColor(0x1A, 0x4A, 0x8C)
PURPLE    = RGBColor(0x5C, 0x2A, 0x9D)
AMBER     = RGBColor(0xB0, 0x70, 0x20)
TEAL      = RGBColor(0x0F, 0x6B, 0x6B)
RED_SOFT  = RGBColor(0xB8, 0x30, 0x20)

CARD_BG   = RGBColor(0x1a, 0x1a, 0x1a)
CARD_BRD  = RGBColor(0x2e, 0x2e, 0x2e)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]   # completely blank layout

def slide():
    s = prs.slides.add_slide(blank)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG_DARK
    return s

def box(s, x, y, w, h, bg=None, border=None, radius=0):
    shape = s.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    fill = shape.fill
    if bg:
        fill.solid(); fill.fore_color.rgb = bg
    else:
        fill.background()
    line = shape.line
    if border:
        line.color.rgb = border; line.width = Pt(0.75)
    else:
        line.fill.background()
    shape.shadow.inherit = False
    # rounded corners via XML adjustment
    if radius:
        sp = shape._element
        spPr = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
        if spPr is not None:
            prstGeom = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
            if prstGeom is not None:
                prstGeom.set('prst', 'roundRect')
                avLst = prstGeom.find('{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
                if avLst is None:
                    avLst = etree.SubElement(prstGeom, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
                gd = etree.SubElement(avLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}gd')
                gd.set('name', 'adj'); gd.set('fmla', 'val %d' % radius)
    return shape

def txt(s, text, x, y, w, h,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = 'Segoe UI'
    return tb

def label(s, text, x, y, w, h, size=10, color=MUTED, bold=True, align=PP_ALIGN.LEFT):
    return txt(s, text, x, y, w, h, size=size, bold=bold, color=color, align=align)

def card(s, x, y, w, h, title, body, icon='', bg=CARD_BG, brd=CARD_BRD,
         title_color=WHITE, body_color=MUTED, icon_size=22, title_size=13, body_size=10):
    box(s, x, y, w, h, bg=bg, border=brd, radius=8000)
    cx = x + Inches(0.18)
    cw = w  - Inches(0.36)
    oy = y  + Inches(0.18)
    if icon:
        txt(s, icon, cx, oy, Inches(0.5), Inches(0.4), size=icon_size, color=WHITE)
        oy += Inches(0.35)
    txt(s, title, cx, oy, cw, Inches(0.3), size=title_size, bold=True, color=title_color)
    txt(s, body,  cx, oy + Inches(0.28), cw, h - Inches(0.9), size=body_size, color=body_color)

def pill_label(s, text, x, y, bg_color, text_color=WHITE, size=9):
    # estimate width
    w = Inches(len(text) * 0.085 + 0.28)
    box(s, x, y, w, Inches(0.28), bg=bg_color, radius=10000)
    txt(s, text, x + Inches(0.1), y + Inches(0.03), w, Inches(0.25),
        size=size, bold=True, color=text_color, align=PP_ALIGN.LEFT)
    return x + w + Inches(0.12)

def divider(s, x, y, w):
    ln = s.shapes.add_connector(1, x, y, x+w, y)
    ln.line.color.rgb = RGBColor(0x2e, 0x2e, 0x2e)
    ln.line.width = Pt(0.5)

def tag_line(s, text):
    txt(s, text, Inches(0.6), Inches(0.38), W - Inches(1.2), Inches(0.3),
        size=9, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

# ── SLIDE 1 ── Title ─────────────────────────────────────────────────────────
s1 = slide()
tag_line(s1, 'FLOURISH  ·  PERSONAL GROWTH APP')

txt(s1, 'Honest patterns.', Inches(1), Inches(1.4), Inches(11.3), Inches(1.1),
    size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s1, 'Better days.', Inches(1), Inches(2.35), Inches(11.3), Inches(1.1),
    size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# subtle subtitle
txt(s1,
    'A personal growth app that surfaces your finance drains, attention leaks, and\n'
    'nature misalignment — automatically from 3-tap daily logs.\n'
    'No account needed. No data on our servers.',
    Inches(2), Inches(3.55), Inches(9.3), Inches(1.1),
    size=13, color=MUTED, align=PP_ALIGN.CENTER)

# three pills
px = Inches(3.4)
py = Inches(4.9)
px = pill_label(s1, '  Free tier — full offline  ', px, py,
                RGBColor(0x1a, 0x0e, 0x38))
px = pill_label(s1, '  Pro — ₹100 one-time  ', px, py,
                RGBColor(0x0b, 0x2e, 0x14))
pill_label(s1,    '  Know yourself in 3 words  ', px, py,
                RGBColor(0x12, 0x20, 0x38))

# bottom tagline
txt(s1, 'flourish.is-a.dev', Inches(1), Inches(6.9), Inches(11.3), Inches(0.35),
    size=10, color=MUTED, align=PP_ALIGN.CENTER)

# ── SLIDE 2 ── Free Features ─────────────────────────────────────────────────
s2 = slide()
tag_line(s2, 'FREE  ·  ALWAYS  ·  NO ACCOUNT NEEDED')

txt(s2, 'Everything you need, offline.', Inches(1), Inches(0.7), Inches(11.3), Inches(0.7),
    size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

feats = [
    ('⚡', '3-tap Quick Log',
     'Activity · Direction · Nature score. Log in under 20 seconds. Finance entry optional.',
     CARD_BG, CARD_BRD),
    ('📊', 'Dashboard + Streaks',
     'Balance radar, habit streaks, category drills. All derived from your own entries.',
     RGBColor(0x09,0x1d,0x0e), RGBColor(0x1A,0x6B,0x3C)),
    ('🔁', 'Insights + Habit Loops',
     'Finance drains, attention leaks, learning gaps surface automatically from your data.',
     CARD_BG, CARD_BRD),
    ('📖', 'Weekly Review',
     'Local review with score, biggest win, watch-list, and one clear focus for next week.',
     RGBColor(0x09,0x12,0x22), RGBColor(0x1A,0x4A,0x8C)),
    ('🌿', 'Nature + Circadian Clock',
     'Live 24-hour clock, nature tickers, Vedic alignment guide. Real-time period awareness.',
     CARD_BG, CARD_BRD),
    ('💾', 'Backup & Restore',
     'One-tap portable code holds all entries. Paste on any device. No account, no expiry.',
     RGBColor(0x1a,0x12,0x05), RGBColor(0xB0,0x70,0x20)),
]

cols = 3
cw = Inches(3.9); ch = Inches(1.55); gap_x = Inches(0.18); gap_y = Inches(0.16)
start_x = Inches(0.6); start_y = Inches(1.65)

for i, (ico, title, body, bg, brd) in enumerate(feats):
    col = i % cols; row = i // cols
    cx = start_x + col * (cw + gap_x)
    cy = start_y + row * (ch + gap_y)
    card(s2, cx, cy, cw, ch, title, body, icon=ico, bg=bg, brd=brd,
         title_size=13, body_size=10)

# ── SLIDE 3 ── Pro Features ───────────────────────────────────────────────────
s3 = slide()
tag_line(s3, 'PRO  ·  ₹100 ONE-TIME  ·  ≈ $1.20')

txt(s3, 'Two features that remove the last friction.', Inches(1), Inches(0.7),
    Inches(11.3), Inches(0.7), size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# big two cards
bw = Inches(5.8); bh = Inches(2.45)
bx1 = Inches(0.7); bx2 = Inches(6.83); by = Inches(1.6)

# Voice Log
box(s3, bx1, by, bw, bh,
    bg=RGBColor(0x14,0x08,0x28), border=RGBColor(0x5C,0x2A,0x9D), radius=8000)
txt(s3, '🎙️', bx1+Inches(0.22), by+Inches(0.2), Inches(0.7), Inches(0.7), size=30, color=WHITE)
txt(s3, 'Voice Logging', bx1+Inches(0.22), by+Inches(0.88),
    bw-Inches(0.44), Inches(0.45), size=19, bold=True, color=WHITE)
txt(s3, 'Speak your entry anywhere in the app.\nTranscribed instantly — no audio stored.\nActive throughout, not just in log entry.',
    bx1+Inches(0.22), by+Inches(1.3), bw-Inches(0.44), Inches(1.1), size=11, color=MUTED)

# Sync
box(s3, bx2, by, bw, bh,
    bg=RGBColor(0x05,0x1a,0x0d), border=RGBColor(0x1A,0x6B,0x3C), radius=8000)
txt(s3, '☁️', bx2+Inches(0.22), by+Inches(0.2), Inches(0.7), Inches(0.7), size=30, color=WHITE)
txt(s3, 'Cross-device Sync', bx2+Inches(0.22), by+Inches(0.88),
    bw-Inches(0.44), Inches(0.45), size=19, bold=True, color=WHITE)
txt(s3, 'Entries sync to your own Google account —\nnot ours. You own your cloud. We just merge.\nSign in once, data follows everywhere.',
    bx2+Inches(0.22), by+Inches(1.3), bw-Inches(0.44), Inches(1.1), size=11, color=MUTED)

# three small cards
sw = Inches(3.9); sh = Inches(1.28)
sx = Inches(0.7); sy = Inches(4.2)
smalls = [
    ('🌙', 'Monthly Reflection',
     'Dark card — narrative sentence, 4 metric tiles, peak day, dominant mode, friction %'),
    ('🧬', 'Rhythm DNA Chart',
     'Day-of-week bar chart. Bar height = avg hrs. Color = directed %. Strongest day highlighted.'),
    ('⭐', 'Early Access',
     'Habit coaching letters, accountability partner, and more as they ship to Pro first.'),
]
for i, (ico, t, b) in enumerate(smalls):
    cx = sx + i * (sw + Inches(0.18))
    card(s3, cx, sy, sw, sh, t, b, icon=ico, title_size=12, body_size=10)

# ── SLIDE 4 ── Payment Flow ───────────────────────────────────────────────────
s4 = slide()
tag_line(s4, 'HOW TO GET PRO  ·  PAY ONCE  ·  UNLOCK FOREVER')

txt(s4, 'Pay ₹100. Verify email. Done.', Inches(1), Inches(0.7),
    Inches(11.3), Inches(0.65), size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# price badge
px_b = Inches(1.0); py_b = Inches(1.65)
bw_b = Inches(3.4); bh_b = Inches(2.4)
box(s4, px_b, py_b, bw_b, bh_b,
    bg=RGBColor(0x0a,0x0a,0x14), border=RGBColor(0x2e,0x2e,0x55), radius=8000)
txt(s4, '₹100', px_b, py_b+Inches(0.45), bw_b, Inches(0.8),
    size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s4, 'one-time  ·  ≈ $1.20', px_b, py_b+Inches(1.25), bw_b, Inches(0.3),
    size=11, color=MUTED, align=PP_ALIGN.CENTER)
txt(s4, 'GPay  ·  PhonePe  ·  Any UPI', px_b, py_b+Inches(1.6), bw_b, Inches(0.3),
    size=11, bold=True, color=RGBColor(0x60,0x60,0xff), align=PP_ALIGN.CENTER)

# consent note
txt(s4, 'One-time · No subscription · No auto-renewal\nDigital access · Non-refundable once code issued',
    px_b, py_b+Inches(1.98), bw_b, Inches(0.45),
    size=9, color=MUTED, align=PP_ALIGN.CENTER)

# steps
steps = [
    ('01', 'Open Settings → Get Pro',
     'Tap the "Get Pro — ₹100" button in Account & Sync'),
    ('02', 'Enter email + agree to terms',
     'Your email receives the activation code. One-time, no renewals.'),
    ('03', 'Pay via GPay / PhonePe / UPI',
     'Razorpay secure checkout. Any UPI method accepted.'),
    ('04', 'Enter OTP from your email',
     '6-digit code sent instantly. Enter it to verify and activate.'),
    ('05', 'Pro active — voice + sync on',
     'Status saved to device. No re-login unless you clear app data.'),
]
sx_s = Inches(4.85); sy_s = Inches(1.55)
sw_s = Inches(7.9)
for i, (num, title, body) in enumerate(steps):
    ey = sy_s + i * Inches(1.03)
    txt(s4, num, sx_s, ey, Inches(0.5), Inches(0.35),
        size=10, bold=True, color=MUTED)
    txt(s4, title, sx_s+Inches(0.5), ey, sw_s-Inches(0.5), Inches(0.32),
        size=13, bold=True, color=WHITE)
    txt(s4, body, sx_s+Inches(0.5), ey+Inches(0.3), sw_s-Inches(0.5), Inches(0.3),
        size=10, color=MUTED)
    if i < len(steps)-1:
        divider(s4, sx_s, ey+Inches(0.72), sw_s)

# ── SLIDE 5 ── Privacy & Know Yourself ───────────────────────────────────────
s5 = slide()
tag_line(s5, 'PRIVACY  ·  YOUR DATA  ·  YOUR CLOUD')

txt(s5, 'Your data never leaves your hands.', Inches(1), Inches(0.7),
    Inches(11.3), Inches(0.65), size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s5,
    'Flourish collects nothing. Stores nothing on its own servers. Processes everything locally.\n'
    'Your entries live on your device — or your own Google account when Pro sync is on.',
    Inches(1.8), Inches(1.5), Inches(9.7), Inches(0.75),
    size=12, color=MUTED, align=PP_ALIGN.CENTER)

prv = [
    ('🔒', 'No server storage'),
    ('📵', 'Works fully offline'),
    ('🧾', 'No ads, no tracking'),
    ('☁️', 'Your cloud = your Google account'),
    ('🎤', 'Audio never recorded or sent'),
    ('💳', 'Payments via Razorpay (PCI-DSS)'),
]
pw = Inches(3.85); ph = Inches(0.88)
pgap = Inches(0.16)
psx = Inches(0.7); psy = Inches(2.5)
for i, (ico, t) in enumerate(prv):
    col = i % 3; row = i // 3
    cx = psx + col*(pw+pgap)
    cy = psy + row*(ph+pgap)
    box(s5, cx, cy, pw, ph,
        bg=RGBColor(0x14,0x14,0x14), border=CARD_BRD, radius=6000)
    txt(s5, ico+' '+t, cx+Inches(0.18), cy+Inches(0.27),
        pw-Inches(0.3), Inches(0.4), size=12, bold=True, color=WHITE)

# closing thought
box(s5, Inches(2.0), Inches(5.55), Inches(9.33), Inches(1.28),
    bg=RGBColor(0x09,0x09,0x16), border=RGBColor(0x2e,0x2e,0x55), radius=8000)
txt(s5, '"Log. Reflect. Shift."',
    Inches(2.0), Inches(5.72), Inches(9.33), Inches(0.55),
    size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s5, 'Three words. Every entry moves you closer to knowing yourself honestly.',
    Inches(2.0), Inches(6.2), Inches(9.33), Inches(0.4),
    size=11, color=MUTED, align=PP_ALIGN.CENTER)

# ── save ─────────────────────────────────────────────────────────────────────
out = '/home/kali/Coding/health-saver-app/flourish-deck.pptx'
prs.save(out)
print('Saved:', out)
